import csv
import io
from pathlib import Path

from PIL import Image
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

import markdown as md

SUPPORTED_EXTENSIONS = {
    ".docx", ".doc", ".xlsx", ".xls", ".csv",
    ".pptx", ".ppt", ".png", ".jpg", ".jpeg",
    ".bmp", ".tiff", ".txt", ".md", ".html",
}


def is_supported(filename: str) -> bool:
    return Path(filename).suffix.lower() in SUPPORTED_EXTENSIONS


def convert_to_pdf(input_path: Path, output_path: Path) -> Path:
    """Convert a file to PDF. Returns the output PDF path."""
    ext = input_path.suffix.lower()

    converters = {
        ".txt": _convert_text,
        ".md": _convert_markdown,
        ".csv": _convert_csv,
        ".png": _convert_image,
        ".jpg": _convert_image,
        ".jpeg": _convert_image,
        ".bmp": _convert_image,
        ".tiff": _convert_image,
        ".docx": _convert_docx,
        ".doc": _convert_docx,
        ".xlsx": _convert_xlsx,
        ".xls": _convert_xlsx,
        ".pptx": _convert_pptx,
        ".ppt": _convert_pptx,
        ".html": _convert_html,
    }

    converter = converters.get(ext)
    if not converter:
        raise ValueError(f"Unsupported file type: {ext}")

    return converter(input_path, output_path)


def _convert_text(input_path: Path, output_path: Path) -> Path:
    """Convert plain text to PDF."""
    text = input_path.read_text(encoding="utf-8", errors="replace")
    styles = getSampleStyleSheet()
    style = ParagraphStyle(
        "Body", parent=styles["Normal"], fontSize=10,
        leading=14, fontName="Courier"
    )

    doc = SimpleDocTemplate(str(output_path), pagesize=letter)
    story = []
    for line in text.split("\n"):
        line = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        line = line.replace(" ", "&nbsp;") if line.strip() == "" else line
        story.append(Paragraph(line or "&nbsp;", style))

    doc.build(story)
    return output_path


def _convert_markdown(input_path: Path, output_path: Path) -> Path:
    """Convert Markdown to PDF via HTML intermediate."""
    text = input_path.read_text(encoding="utf-8", errors="replace")
    html_content = md.markdown(text, extensions=["tables", "fenced_code"])
    return _html_to_pdf_reportlab(html_content, output_path)


def _convert_html(input_path: Path, output_path: Path) -> Path:
    """Convert HTML file to PDF."""
    html_content = input_path.read_text(encoding="utf-8", errors="replace")
    return _html_to_pdf_reportlab(html_content, output_path)


def _html_to_pdf_reportlab(html_content: str, output_path: Path) -> Path:
    """Render HTML content to PDF using reportlab (basic)."""
    # Strip tags for a simple text-based rendering
    import re
    clean = re.sub(r"<[^>]+>", " ", html_content)
    clean = clean.replace("&nbsp;", " ").replace("&amp;", "&")
    clean = clean.replace("&lt;", "<").replace("&gt;", ">")

    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(str(output_path), pagesize=letter)
    story = []
    for line in clean.split("\n"):
        line = line.strip()
        if line:
            safe = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, styles["Normal"]))
            story.append(Spacer(1, 6))

    if not story:
        story.append(Paragraph("(empty document)", styles["Normal"]))

    doc.build(story)
    return output_path


def _convert_image(input_path: Path, output_path: Path) -> Path:
    """Convert an image file to a PDF page."""
    img = Image.open(input_path)
    if img.mode in ("RGBA", "P"):
        img = img.convert("RGB")

    img_width, img_height = img.size
    page_width, page_height = letter
    margin = 0.5 * inch

    available_w = page_width - 2 * margin
    available_h = page_height - 2 * margin

    scale = min(available_w / img_width, available_h / img_height, 1.0)
    draw_w = img_width * scale
    draw_h = img_height * scale

    c = canvas.Canvas(str(output_path), pagesize=letter)
    x = (page_width - draw_w) / 2
    y = (page_height - draw_h) / 2
    c.drawImage(ImageReader(img), x, y, draw_w, draw_h)
    c.save()
    return output_path


def _convert_csv(input_path: Path, output_path: Path) -> Path:
    """Convert CSV to a table-based PDF."""
    with open(input_path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        data = [row for row in reader]

    if not data:
        data = [["(empty file)"]]

    # Truncate very wide tables
    max_cols = 10
    data = [row[:max_cols] for row in data]

    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(str(output_path), pagesize=A4)
    story = []

    # Style the table
    table_style = TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#334155")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f1f5f9")]),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
    ])

    # Split into chunks of 50 rows per table to fit pages
    chunk_size = 50
    for i in range(0, len(data), chunk_size):
        chunk = data[i:i + chunk_size]
        t = Table(chunk)
        t.setStyle(table_style)
        story.append(t)
        story.append(Spacer(1, 12))

    doc.build(story)
    return output_path


def _convert_docx(input_path: Path, output_path: Path) -> Path:
    """Convert DOCX to PDF using python-docx + reportlab."""
    from docx import Document

    doc = Document(str(input_path))
    styles = getSampleStyleSheet()

    pdf = SimpleDocTemplate(str(output_path), pagesize=letter)
    story = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            story.append(Spacer(1, 6))
            continue

        safe = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

        style_name = para.style.name if para.style and para.style.name else ""
        if style_name.startswith("Heading"):
            level = style_name[-1] if style_name[-1].isdigit() else "1"
            style_name = f"Heading{min(int(level), 6)}"
            if style_name in styles:
                story.append(Paragraph(safe, styles[style_name]))
            else:
                story.append(Paragraph(f"<b>{safe}</b>", styles["Normal"]))
        else:
            story.append(Paragraph(safe, styles["Normal"]))

    if not story:
        story.append(Paragraph("(empty document)", styles["Normal"]))

    pdf.build(story)
    return output_path


def _convert_xlsx(input_path: Path, output_path: Path) -> Path:
    """Convert XLSX to PDF table."""
    from openpyxl import load_workbook

    wb = load_workbook(str(input_path), read_only=True, data_only=True)
    styles = getSampleStyleSheet()
    pdf = SimpleDocTemplate(str(output_path), pagesize=A4)
    story = []

    for sheet in wb.worksheets:
        story.append(Paragraph(f"<b>Sheet: {sheet.title}</b>", styles["Heading2"]))
        story.append(Spacer(1, 8))

        data = []
        for row in sheet.iter_rows(max_row=200, max_col=10, values_only=True):
            data.append([str(cell) if cell is not None else "" for cell in row])

        if data:
            table_style = TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#334155")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f1f5f9")]),
            ])
            t = Table(data)
            t.setStyle(table_style)
            story.append(t)
        else:
            story.append(Paragraph("(empty sheet)", styles["Normal"]))

        story.append(Spacer(1, 20))

    wb.close()
    if not story:
        story.append(Paragraph("(empty workbook)", styles["Normal"]))

    pdf.build(story)
    return output_path


def _convert_pptx(input_path: Path, output_path: Path) -> Path:
    """Convert PPTX to PDF — one page per slide with text content."""
    from pptx import Presentation

    prs = Presentation(str(input_path))
    styles = getSampleStyleSheet()
    pdf = SimpleDocTemplate(str(output_path), pagesize=letter)
    story = []

    for i, slide in enumerate(prs.slides, 1):
        story.append(Paragraph(f"<b>Slide {i}</b>", styles["Heading2"]))
        story.append(Spacer(1, 8))

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        safe = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                        story.append(Paragraph(safe, styles["Normal"]))
                        story.append(Spacer(1, 4))

        story.append(Spacer(1, 20))

    if not story:
        story.append(Paragraph("(empty presentation)", styles["Normal"]))

    pdf.build(story)
    return output_path
